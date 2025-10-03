# audit_lib.py — pure library (NO streamlit import)
from __future__ import annotations

import re
import pathlib
from collections import Counter, defaultdict
from typing import List, Tuple, Dict, Optional

# Lazy imports to keep deps optional until needed
docx = None
Presentation = None
pd = None
etree = None

def _lazy_imports():
    global docx, Presentation, pd, etree
    if docx is None:
        try:
            import docx  # python-docx
        except Exception:
            docx = None
    if Presentation is None:
        try:
            from pptx import Presentation  # python-pptx
        except Exception:
            Presentation = None
    if pd is None:
        try:
            import pandas as pd
        except Exception:
            pd = None
    if etree is None:
        try:
            from lxml import etree  # robust XML handling for XLIFF/MXLIFF
        except Exception:
            etree = None

# ---- Numeric normalization helpers ----
ARABIC_INDIC = dict(zip("٠١٢٣٤٥٦٧٨٩", "0123456789"))
EASTERN_ARABIC_INDIC = dict(zip("۰۱۲۳۴۵۶۷۸۹", "0123456789"))

def normalize_digits(s: str) -> str:
    if not s:
        return s
    s = s.translate(str.maketrans(ARABIC_INDIC))
    s = s.translate(str.maketrans(EASTERN_ARABIC_INDIC))
    return s

def normalize_number_str(num: str) -> str:
    # Normalize 1,234.56 / 1.234,56 / 1 234,56 / ١٢٣,٤٥ -> 1234.56
    num = normalize_digits(num)
    num = re.sub(r"[\u00A0\u2000-\u200A\u202F\u205F\u3000]", " ", num)

    if "," in num and "." in num:
        if num.rfind(",") > num.rfind("."):
            num = num.replace(".", "")
            num = num.replace(",", ".")
        else:
            num = num.replace(",", "")
    else:
        if "," in num:
            parts = num.split(",")
            if len(parts) == 2 and len(parts[1]) in (1, 2):
                num = num.replace(",", ".")
            else:
                num = num.replace(",", "")
        if "." in num and num.count(".") > 1:
            last = num.rfind(".")
            if len(num) - last - 1 in (1, 2):
                num = num[:last].replace(".", "") + "." + num[last+1:]
            else:
                num = num.replace(".", "")

    num = re.sub(r"\s(?=\d{3}(\D|$))", "", num)
    return num

def as_float_safe(num: str):
    try:
        return float(normalize_number_str(num))
    except Exception:
        return None

# ---- Patterns ----
UNIT_REGEX = r"(?:%|％|°[CF]?|℃|℉|V|mV|kV|A|mA|µA|Ω|ohm|W|kW|MW|Wh|kWh|mAh|J|kJ|Nm|N·m|N-m|Pa|kPa|MPa|bar|mbar|mmHg|psi|Hz|kHz|MHz|GHz|rpm|r/min|s|min|h|hr|d|ms|µs|mol|ppm|ppb|pH|lx|dB|m|cm|mm|km|µm|nm|in|ft|yd|L|mL|µL|kg|g|mg|µg|lb|oz|°|deg|C|F|m/s|km/h|IU|UI/mL|CFU/g)"
NUM_REGEX  = r"[+\-]?(?:\d|\u0660-\u0669|\u06F0-\u06F9)[\d\u0660-\u0669\u06F0-\u06F9\s.,]*\d"

PAIR_REGEX = re.compile(
    rf"(?P<num>{NUM_REGEX})\s*(?P<unit>{UNIT_REGEX})|(?P<unit2>{UNIT_REGEX})\s*(?P<num2>{NUM_REGEX})"
)
PURE_NUM_REGEX = re.compile(
    rf"(?<![A-Za-z\u0600-\u06FF\u4E00-\u9FFF])(?P<num>{NUM_REGEX})(?![A-Za-z\u0600-\u06FF\u4E00-\u9FFF])"
)

def normalize_unit(u: str) -> str:
    u = u.replace(" ", "")
    u = u.replace("ohm", "Ω")
    u = u.replace("r/min", "rpm")
    u = u.replace("deg", "°")
    u = u.replace("N-m", "N·m")
    return u

def canonical_pair(num: str, unit: str):
    n_str = normalize_number_str(num)
    f = as_float_safe(n_str)
    return (n_str, normalize_unit(unit), f if f is not None else float("nan"))

def canonical_number(num: str):
    n_str = normalize_number_str(num)
    f = as_float_safe(n_str)
    return (n_str, f if f is not None else float("nan"))

def find_pairs(text: str):
    out = []
    for m in PAIR_REGEX.finditer(text):
        if m.group("num") and m.group("unit"):
            num, unit = m.group("num"), m.group("unit")
            s, e = m.span()
        else:
            num, unit = m.group("num2"), m.group("unit2")
            s, e = m.span()
        out.append((num, unit, s, e))
    return out

def find_pure_numbers(text: str):
    return [(m.group("num"), m.start(), m.end()) for m in PURE_NUM_REGEX.finditer(text)]

# ---- Readers (txt, docx, pptx, xlsx/csv/tsv, XLIFF/MXLIFF) ----
def _read_text_txt(path: pathlib.Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

def _read_text_docx(path: pathlib.Path) -> str:
    _lazy_imports()
    if docx is None:
        raise RuntimeError("python-docx not installed. pip install python-docx")
    d = docx.Document(str(path))
    parts = []
    for p in d.paragraphs:
        parts.append(p.text)
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)

def _read_text_pptx(path: pathlib.Path) -> str:
    _lazy_imports()
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. pip install python-pptx")
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)

def _read_text_tabular(path: pathlib.Path) -> str:
    _lazy_imports()
    if pd is None:
        raise RuntimeError("pandas not installed. pip install pandas openpyxl")
    ext = path.suffix.lower()
    if ext == ".xlsx":
        df_dict = pd.read_excel(str(path), sheet_name=None, dtype=str)
        parts = []
        for name, df in df_dict.items():
            parts.append(f"### SHEET: {name} ###")
            parts.append(df.fillna("").to_string(index=False))
        return "\n".join(parts)
    elif ext == ".csv":
        df = pd.read_csv(str(path), dtype=str, encoding="utf-8", keep_default_na=False, na_values=[])
        return df.fillna("").to_string(index=False)
    else:  # .tsv
        df = pd.read_csv(str(path), dtype=str, sep="\t", encoding="utf-8", keep_default_na=False, na_values=[])
        return df.fillna("").to_string(index=False)

def read_text_from_file(path: pathlib.Path) -> str:
    ext = path.suffix.lower()
    if ext in [".txt", ".md"]:
        return _read_text_txt(path)
    if ext == ".docx":
        return _read_text_docx(path)
    if ext == ".pptx":
        return _read_text_pptx(path)
    if ext in [".xlsx", ".csv", ".tsv"]:
        return _read_text_tabular(path)
    if ext in [".xliff", ".xlf", ".sdlxliff", ".mxliff"]:
        raise RuntimeError("For XLIFF/MXLIFF use read_xliff_pair(path) and pass source+target separately.")
    raise RuntimeError(f"Unsupported file type: {ext}")

# ---- XLIFF/MXLIFF ----
def read_xliff_pair(path: pathlib.Path) -> tuple[str, str]:
    _lazy_imports()
    if etree is None:
        raise RuntimeError("lxml not installed. pip install lxml")
    parser = etree.XMLParser(recover=True, huge_tree=True)
    root = etree.parse(str(path), parser).getroot()
    sources = root.xpath('.//*[local-name()="source"]')
    targets = root.xpath('.//*[local-name()="target"]')

    def join_text(nodes):
        parts = []
        for n in nodes:
            parts.append("".join(n.itertext()))
        return "\n".join(parts)

    src_text = join_text(sources)
    tgt_text = join_text(targets)

    if not src_text:
        seg_sources = root.xpath('.//*[local-name()="seg-source"]')
        if seg_sources:
            src_text = "\n".join("".join(n.itertext()) for n in seg_sources)
    return src_text, tgt_text

# ---- Comparisons ----
def _compare_pairs(src_text: str, tgt_text: str):
    src_text_norm = normalize_digits(src_text)
    tgt_text_norm = normalize_digits(tgt_text)
    src_pairs = [(m[0], m[1], m[2], m[3]) for m in find_pairs(src_text_norm)]
    tgt_pairs = [(m[0], m[1], m[2], m[3]) for m in find_pairs(tgt_text_norm)]

    src_keyed = defaultdict(list)
    for num, unit, s, e in src_pairs:
        cn, cu, cf = canonical_pair(num, unit)
        src_keyed[cu].append((cn, cf, s, e, num, unit))
    tgt_keyed = defaultdict(list)
    for num, unit, s, e in tgt_pairs:
        cn, cu, cf = canonical_pair(num, unit)
        tgt_keyed[cu].append((cn, cf, s, e, num, unit))

    missing, extra, changed = [], [], []
    units = set(list(src_keyed.keys()) + list(tgt_keyed.keys()))
    for unit in units:
        src_list = sorted([x[0] for x in src_keyed[unit]])
        tgt_list = sorted([x[0] for x in tgt_keyed[unit]])
        src_count = Counter(src_list)
        tgt_count = Counter(tgt_list)

        for k in (src_count - tgt_count).elements():
            cn, cf, s, e, onum, ounit = next(x for x in src_keyed[unit] if x[0] == k)
            missing.append({"unit": unit, "value": k, "orig_num": onum, "orig_unit": ounit, "note": "Present in source, missing in target"})
        for k in (tgt_count - src_count).elements():
            cn, cf, s, e, onum, ounit = next(x for x in tgt_keyed[unit] if x[0] == k)
            extra.append({"unit": unit, "value": k, "orig_num": onum, "orig_unit": ounit, "note": "Present in target, not in source"})

        m = min(len(src_list), len(tgt_list))
        for i in range(m):
            s_val, t_val = src_list[i], tgt_list[i]
            if s_val != t_val:
                sf, tf = as_float_safe(s_val), as_float_safe(t_val)
                if sf is None or tf is None or abs(sf - tf) > 1e-9:
                    changed.append({"unit": unit, "source_value": s_val, "target_value": t_val, "note": "Numeric value differs"})
    return missing, extra, changed

def _compare_pure_numbers(src_text: str, tgt_text: str):
    src_vals = [canonical_number(n)[0] for n, _, _ in find_pure_numbers(normalize_digits(src_text))]
    tgt_vals = [canonical_number(n)[0] for n, _, _ in find_pure_numbers(normalize_digits(tgt_text))]
    return Counter(src_vals), Counter(tgt_vals)

# ---- Public API ----
def audit_files(source_path, target_path, out_xlsx: Optional[str | pathlib.Path] = None):
    _lazy_imports()
    src_p = pathlib.Path(source_path)
    tgt_p = pathlib.Path(target_path)

    # Build source/target text depending on types
    if src_p.suffix.lower() in [".xliff", ".xlf", ".sdlxliff", ".mxliff"]:
        src_source_text, _ = read_xliff_pair(src_p)
        if tgt_p.suffix.lower() in [".xliff", ".xlf", ".sdlxliff", ".mxliff"]:
            _, tgt_target_text = read_xliff_pair(tgt_p)
            source_text = src_source_text
            target_text = tgt_target_text
        else:
            source_text = src_source_text
            target_text = read_text_from_file(tgt_p)
    else:
        source_text = read_text_from_file(src_p)
        if tgt_p.suffix.lower() in [".xliff", ".xlf", ".sdlxliff", ".mxliff"]:
            _, target_text = read_xliff_pair(tgt_p)
        else:
            target_text = read_text_from_file(tgt_p)

    # Compare
    missing_pairs, extra_pairs, changed_pairs = _compare_pairs(source_text, target_text)
    src_num_counts, tgt_num_counts = _compare_pure_numbers(source_text, target_text)

    if pd is None:
        raise RuntimeError("pandas not installed. pip install pandas openpyxl")

    # Build DataFrames
    def mkdf(rows, cols): return pd.DataFrame(rows) if rows else pd.DataFrame(columns=cols)

    df_missing_pairs = mkdf(
        [{"Type":"MissingPair","Unit":i["unit"],"Value":i["value"],"OriginalNumber":i["orig_num"],"OriginalUnit":i["orig_unit"],"Note":i["note"]} for i in missing_pairs],
        ["Type","Unit","Value","OriginalNumber","OriginalUnit","Note"]
    )
    df_extra_pairs = mkdf(
        [{"Type":"ExtraPair","Unit":i["unit"],"Value":i["value"],"OriginalNumber":i["orig_num"],"OriginalUnit":i["orig_unit"],"Note":i["note"]} for i in extra_pairs],
        ["Type","Unit","Value","OriginalNumber","OriginalUnit","Note"]
    )
    df_changed_pairs = mkdf(
        [{"Unit":i["unit"],"SourceValue":i["source_value"],"TargetValue":i["target_value"],"Note":i["note"]} for i in changed_pairs],
        ["Unit","SourceValue","TargetValue","Note"]
    )

    diff_missing = (src_num_counts - tgt_num_counts).items()
    diff_extra   = (tgt_num_counts - src_num_counts).items()
    df_pure_missing = mkdf(
        [{"Type":"PureNumber","Value":v,"CountDiff":c,"Note":"More in source than target"} for v,c in diff_missing],
        ["Type","Value","CountDiff","Note"]
    )
    df_pure_extra = mkdf(
        [{"Type":"PureNumber","Value":v,"CountDiff":c,"Note":"More in target than source"} for v,c in diff_extra],
        ["Type","Value","CountDiff","Note"]
    )

    df_summary = pd.DataFrame([{
        "MissingPairs": len(df_missing_pairs),
        "ExtraPairs": len(df_extra_pairs),
        "ChangedPairs": len(df_changed_pairs),
        "PureNumbers_MissingKinds": len(df_pure_missing),
        "PureNumbers_ExtraKinds": len(df_pure_extra),
    }])

    result = {
        "Summary": df_summary,
        "Missing_in_Target": df_missing_pairs,
        "Extra_in_Target": df_extra_pairs,
        "Value_Changed": df_changed_pairs,
        "PureNums_Missing": df_pure_missing,
        "PureNums_Extra": df_pure_extra,
    }

    if out_xlsx:
        with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
            for name, df in result.items():
                df.to_excel(writer, sheet_name=name[:31], index=False)

    return result

# Optional CLI for quick local testing (no Streamlit)
if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("source")
    ap.add_argument("target")
    ap.add_argument("--out", default="numbers_units_audit.xlsx")
    args = ap.parse_args()
    audit_files(args.source, args.target, out_xlsx=args.out)
    print(f"Done. Wrote {args.out}")
